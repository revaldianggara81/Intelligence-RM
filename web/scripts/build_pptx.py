"""
Generate PowerPoint documentation for Bank Danamon Intelligence RM Platform.
Includes screenshots of all modules with descriptions.
Run: python scripts/build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ── Paths ────────────────────────────────────────────────────────────────────
BASE_DIR   = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SS_DIR     = os.path.join(BASE_DIR, 'docs', 'screenshots')
OUT_PATH   = os.path.join(BASE_DIR, 'docs', 'Danamon_RM_Platform_Documentation.pptx')
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── Theme colours ────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x12, 0x1F)   # dark navy
C_CARD     = RGBColor(0x13, 0x1C, 0x30)   # card bg
C_ACCENT   = RGBColor(0xFF, 0x4E, 0x00)   # Oracle red-orange
C_CYAN     = RGBColor(0x00, 0xCC, 0xFF)   # cyan accent
C_GOLD     = RGBColor(0xFF, 0xB8, 0x30)   # gold/warning
C_GREEN    = RGBColor(0x00, 0xD4, 0x7E)   # success green
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xB0, 0xBE, 0xD8)   # secondary text
C_DARK_TXT = RGBColor(0x1A, 0x1A, 0x2E)
C_DIVIDER  = RGBColor(0x1E, 0x2D, 0x4A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── Slide content definitions ─────────────────────────────────────────────────
SLIDES = [
    # ── COVER ────────────────────────────────────────────────────────────────
    {
        'type': 'cover',
        'title': 'Intelligence RM Platform',
        'subtitle': 'Bank Danamon POC · 2026',
        'tagline': 'Oracle AI · Oracle Database 26ai · OCI Generative AI · Private Agent Factory',
        'badges': ['Oracle ADB 26ai', 'Cohere Command R+', 'Oracle Select AI', 'RAG · Vector Search'],
    },

    # ── AGENDA ───────────────────────────────────────────────────────────────
    {
        'type': 'agenda',
        'title': 'Application Modules',
        'items': [
            ('01', 'Login & Authentication',       'Secure access with role-based permissions'),
            ('02', 'RM Dashboard',                 'Unified RM workspace with KPI overview'),
            ('03', 'Customer 360',                 'Full customer intelligence & AI-powered profiling'),
            ('04', 'Calendar & Appointments',      'RM schedule and customer meeting management'),
            ('05', 'Maturity Reminder',            'AI-driven product maturity analysis & alerts'),
            ('06', 'Product Recommendations',      'Personalized next-product recommendations'),
            ('07', 'Campaign Management',          'Targeted campaign scan & customer segmentation'),
            ('08', 'Portfolio Alerts',             'Real-time risk detection & action triggers'),
            ('09', 'AI Copilot',                   'Conversational AI assistant for RMs'),
            ('10', 'Executive Dashboard',          'Management intelligence & AUM analytics'),
            ('11', 'Compliance',                   'Regulatory monitoring & audit trail'),
            ('12', 'Admin Console',                'User management & system configuration'),
        ],
    },

    # ── 01 LOGIN ─────────────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '01',
        'module_name': 'Login & Authentication',
        'screenshot': '01_login.png',
        'description': (
            'Secure authentication portal for Bank Danamon RM staff. '
            'Supports role-based access: Relationship Manager (RM) and Manager roles, '
            'each with differentiated permissions and dashboard views.'
        ),
        'features': [
            'JWT-based session authentication with 8-hour token expiry',
            'Role-based access control: RM (anisa, budi, dewi) and Manager roles',
            'Oracle ADB backend — credentials validated against USERS table',
            'Oracle Private Agent Factory (PAF) integration badge displayed on login',
            'Demo credential hints for POC environment',
        ],
    },

    # ── 02 DASHBOARD ─────────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '02',
        'module_name': 'RM Dashboard',
        'screenshot': '02_dashboard.png',
        'description': (
            'Central workspace for Relationship Managers. Provides a real-time snapshot '
            'of portfolio health, upcoming actions, and quick access to all platform modules. '
            'Data is pulled live from Oracle ADB 26ai via REST API.'
        ),
        'features': [
            'Portfolio KPI widgets: Total AUM, active customers, alerts count, pending renewals',
            'Quick-action shortcuts to Maturity Reminder, Product Reco, Campaigns, Alerts',
            'Top customers by AUM with risk profile and KYC status badges',
            'Upcoming appointments calendar strip with customer context',
            'Notification bell with real-time alert count from Oracle DB',
        ],
    },
    {
        'type': 'module_b',
        'module_no': '02',
        'module_name': 'RM Dashboard — Portfolio Overview',
        'screenshot': '02b_dashboard_scroll.png',
        'description': (
            'Scrolled view of the RM Dashboard showing portfolio breakdown, '
            'customer segment distribution, and recent activity feed.'
        ),
        'features': [
            'AUM breakdown by customer segment (Priority, Privilege, Regular)',
            'Product holding distribution across investment categories',
            'Recent customer interactions and activity log',
            'Expiring product alerts with days-to-maturity countdown',
            'RM performance metric vs. team average comparison',
        ],
    },

    # ── 03 CUSTOMER 360 ───────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '03',
        'module_name': 'Customer 360 — Customer List',
        'screenshot': '03_customer360_list.png',
        'description': (
            'Complete customer management view listing all RM-assigned customers with '
            'segment classification, AUM, KYC status, and alert indicators. '
            'Supports search and filter by segment (All / Prioritas / Privilege / Regular).'
        ),
        'features': [
            'Live customer list from Oracle ADB with real-time KYC & alert status',
            'Segment filter tabs: All, Prioritas, Privilege, Regular',
            'Search bar with instant lookup by name or customer ID',
            'Color-coded segment badges: Prioritas (red), Privilege (blue), Regular (grey)',
            'Alert indicators on customer cards with severity count',
        ],
    },
    {
        'type': 'module',
        'module_no': '03',
        'module_name': 'Customer 360 — Profile & AI Summary',
        'screenshot': '03b_customer360_profile.png',
        'description': (
            'Full 360° customer intelligence view. Displays complete profile with AUM, '
            'risk profile, product holdings, active alerts, KYC status, and an '
            'AI-generated call center interaction summary powered by OCI Cohere.'
        ),
        'features': [
            'Customer header: AUM, risk profile, KYC badge, segment, age',
            'AI Summary — Call Center: auto-generated from interaction transcripts via RAG',
            'Sentiment analysis with interaction quality indicators (positive/negative/neutral)',
            'Key complaint & insight highlights extracted by Oracle Select AI',
            'Recommended follow-up actions ranked by urgency',
        ],
    },
    {
        'type': 'module',
        'module_no': '03',
        'module_name': 'Customer 360 — AI Market Intelligence',
        'screenshot': '03d_customer360_market_intel.png',
        'description': (
            'Personalized AI-generated market intelligence card, tailored to each customer\'s '
            'risk profile and active portfolio. Generated via OCI Cohere Command R+ using '
            'live Indonesia macro-economic context. Results cached 24h in Oracle ADB.'
        ),
        'features': [
            'Live market data strip: USD/IDR, BI Rate, IHSG, Inflasi, ORI028, Emas',
            'Up/down/warn chips with directional color-coding (green/red/amber)',
            '3 personalized investment insight points by AI (max 130 words, Bahasa Indonesia)',
            'Customer-specific context: RISK_PROFILE, TOTAL_AUM, active holdings, alerts',
            '24-hour cache in CUSTOMER_AI_INSIGHTS table with force-refresh option',
        ],
    },
    {
        'type': 'module',
        'module_no': '03',
        'module_name': 'Customer 360 — Sub-modules',
        'screenshot': '03e_customer360_submodules.png',
        'description': (
            'Collapsible sub-modules within Customer 360 providing deep operational detail: '
            'product forecasts, RM meeting notes, and call center transcripts.'
        ),
        'features': [
            'Forecast Current Product: Q1–Q4 projected accumulation per product category',
            'Catatan RM: personal and meeting notes with category filter and add-new form',
            'Transcript Call Center: AI-analyzed interaction logs with timestamps',
            'All sub-modules collapsible/expandable for focused workflow',
            'Vector-embedded customer data enables semantic RAG responses',
        ],
    },

    # ── 04 CALENDAR ──────────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '04',
        'module_name': 'Calendar & Appointments',
        'screenshot': '04_calendar.png',
        'description': (
            'Visual calendar for managing RM customer appointments. Displays all scheduled '
            'meetings with customer context, appointment type, and priority indicators. '
            'Integrated with Oracle ADB RM_APPOINTMENTS table.'
        ),
        'features': [
            'Monthly calendar view with appointment density indicators per day',
            'Appointment list with customer name, type (Meeting/Call/Review), and time',
            'Priority color-coding: high (red), medium (amber), normal (blue)',
            'Customer AUM and segment shown alongside each appointment',
            'Navigation between months with appointment count summary',
        ],
    },

    # ── 05 MATURITY REMINDER ─────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '05',
        'module_name': 'Maturity Reminder — Analysis Input',
        'screenshot': '05_maturity_reminder.png',
        'description': (
            'AI Agent for detecting upcoming product maturities in the RM\'s portfolio. '
            'The RM initiates an analysis run that queries Oracle ADB for expiring products '
            'and generates personalized reinvestment recommendations via OCI Generative AI.'
        ),
        'features': [
            'Scans all active holdings with maturity dates within configurable horizon (30/60/90 days)',
            'AI analysis powered by Oracle Private Agent Factory (PAF) tools',
            'Customer segmentation with AUM and product type context',
            'One-click "Run Analysis" triggers Oracle AI Agent pipeline',
            'History button to review previous analysis runs from Oracle DB',
        ],
    },
    {
        'type': 'module',
        'module_no': '05',
        'module_name': 'Maturity Reminder — AI Analysis Result',
        'screenshot': '05b_maturity_result.png',
        'description': (
            'AI-generated maturity analysis output. The agent identifies customers with '
            'products approaching maturity and produces tailored reinvestment action plans '
            'with product recommendations based on risk profile and current market rates.'
        ),
        'features': [
            'Per-customer maturity summary with product name, amount, and days remaining',
            'AI-recommended reinvestment products matched to customer risk profile',
            'Estimated yield improvement from rolling over to suggested alternatives',
            'Action buttons: Share via Email / WhatsApp, Action Plan Follow-up, Customer 360',
            'Copy-to-clipboard for quick use in CRM or email drafting',
        ],
    },

    # ── 06 PRODUCT RECO ──────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '06',
        'module_name': 'Product Recommendations',
        'screenshot': '06_product_reco.png',
        'description': (
            'AI-powered Next-Best-Product recommendation engine. Analyzes each customer\'s '
            'portfolio gaps, transaction history, risk profile, and market conditions to '
            'surface the most relevant product offers for the RM to present.'
        ),
        'features': [
            'Runs via Oracle PAF Agent with access to PRODUCT_CATALOG and customer holdings',
            'Multi-factor scoring: portfolio gap, risk alignment, AUM potential, urgency',
            'Top-N recommendations per customer with rationale explanation in Bahasa Indonesia',
            'Products ranked by cross-sell probability and estimated revenue contribution',
            'Configurable customer scope: priority segment or full portfolio',
        ],
    },
    {
        'type': 'module',
        'module_no': '06',
        'module_name': 'Product Recommendations — Results',
        'screenshot': '06b_product_reco_result.png',
        'description': (
            'Recommendation analysis results showing per-customer product suggestions '
            'with actionable next steps, share options, and direct link to Customer 360.'
        ),
        'features': [
            'Ranked product list per customer with fit score and key rationale',
            'Market context used in recommendation (BI Rate, IHSG, currency rates)',
            'Estimated AUM uplift and product category diversification impact',
            'Direct CTA to open Customer 360 for immediate consultation preparation',
            'Share via Email / WhatsApp for remote customer outreach',
        ],
    },

    # ── 07 CAMPAIGN MANAGEMENT ────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '07',
        'module_name': 'Campaign Management',
        'screenshot': '07_campaign_mgmt.png',
        'description': (
            'Intelligent campaign targeting tool. Scans the RM\'s entire customer portfolio '
            'to identify the best-fit customers for a specific bank product campaign, '
            'ranked by propensity score using Oracle AI analytics.'
        ),
        'features': [
            'Campaign selection: choose active bank campaigns from Oracle DB catalog',
            'AI scan maps campaign criteria against all customer profiles automatically',
            'Output: prioritized customer shortlist with fit rationale per customer',
            'Segment and AUM filter for targeted campaign execution',
            'Campaign history tracking stored in Oracle ADB for audit and follow-up',
        ],
    },
    {
        'type': 'module',
        'module_no': '07',
        'module_name': 'Campaign Management — Scan Results',
        'screenshot': '07b_campaign_scan_result.png',
        'description': (
            'Campaign scan result showing the ranked list of customers who best match '
            'the selected campaign, with AI-generated talking points for each customer.'
        ),
        'features': [
            'Ranked customer targets with propensity score and match rationale',
            'AI-generated personalized pitch points per customer in Bahasa Indonesia',
            'Customer AUM, risk profile, and existing product context displayed inline',
            'Share result via Email / WhatsApp for offline CRM integration',
            'Action Plan Follow-up to log campaign contact attempts in Oracle DB',
        ],
    },

    # ── 08 PORTFOLIO ALERTS ───────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '08',
        'module_name': 'Portfolio Alerts',
        'screenshot': '08_portfolio_alerts.png',
        'description': (
            'Real-time portfolio risk monitoring and alert management. Detects anomalies '
            'across the RM\'s customer portfolio including AUM decline, idle funds, '
            'KYC expiry, and churn risk signals, then surfaces actionable interventions.'
        ),
        'features': [
            'Alert categories: Churn Risk, Idle Funds, KYC Expiry, Maturity Approaching, AUM Drop',
            'Severity levels: Critical (red), High (orange), Medium (amber), Low (blue)',
            'Auto-detect scan via Oracle AI analyzing all portfolio holdings',
            'Per-alert AI-recommended action with urgency timeline',
            'Alert count badge in sidebar and notification bell for real-time awareness',
        ],
    },
    {
        'type': 'module',
        'module_no': '08',
        'module_name': 'Portfolio Alerts — Detail View',
        'screenshot': '08b_portfolio_alerts_detail.png',
        'description': (
            'Detailed alert view showing customer-level risk breakdown with AI analysis '
            'explaining the root cause of each alert and recommended remediation steps.'
        ),
        'features': [
            'Per-customer alert cards with product name, risk type, and severity badge',
            'AI-generated explanation of the alert trigger in plain Bahasa Indonesia',
            'Recommended action steps ranked by impact and feasibility',
            'One-click access to Customer 360 for immediate deep-dive investigation',
            'Bulk alert dismiss or mark-as-handled for workflow efficiency',
        ],
    },

    # ── 09 AI COPILOT ─────────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '09',
        'module_name': 'AI Copilot',
        'screenshot': '09_ai_copilot.png',
        'description': (
            'Conversational AI assistant for Relationship Managers. Powered by Oracle '
            'Private Agent Factory (PAF) with access to customer data, product catalog, '
            'market intelligence, and bank policies via RAG on Oracle Vector Search.'
        ),
        'features': [
            'Natural language interface — ask about any customer, product, or market condition',
            'Context-aware: Copilot knows the RM\'s portfolio and can reference specific customers',
            'RAG-powered: retrieves relevant data from Oracle ADB vector embeddings',
            'Supports Bahasa Indonesia and English queries seamlessly',
            'Conversation history stored in Oracle DB for continuity across sessions',
        ],
    },
    {
        'type': 'module',
        'module_no': '09',
        'module_name': 'AI Copilot — Query Interface',
        'screenshot': '09b_ai_copilot_input.png',
        'description': (
            'Copilot query input showing the conversational interface with suggested '
            'prompts and response area. The Copilot uses Oracle Select AI and PAF tools '
            'to synthesize answers from live Oracle ADB data.'
        ),
        'features': [
            'Suggested query chips for common RM workflows (portfolio summary, top risks, etc.)',
            'Free-text input with Enter-to-send and send button',
            'Streaming response display with markdown formatting support',
            'Source citation: responses indicate which Oracle DB tables were queried',
            'Clear conversation / history navigation for session management',
        ],
    },

    # ── 10 EXECUTIVE DASHBOARD ────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '10',
        'module_name': 'Executive Dashboard',
        'screenshot': '10_executive_dashboard.png',
        'description': (
            'Management-level intelligence dashboard providing a consolidated view of '
            'portfolio performance, RM productivity, customer health distribution, and '
            'AUM trend analytics. Powered by Oracle ADB 26ai with AI-generated insights.'
        ),
        'features': [
            'AUM Velocity widget: 6-month trend chart with 2-month AI forecast',
            'Customer Health Distribution: pie chart of health scores across portfolio',
            'Customer Intelligence Radar: churn risk vs. opportunity scoring bubble chart',
            'RM Productivity Cockpit: per-RM AUM, customer count, and activity metrics',
            'AI-generated executive insights per widget via OCI Cohere',
        ],
    },
    {
        'type': 'module',
        'module_no': '10',
        'module_name': 'Executive Dashboard — Intelligence Panels',
        'screenshot': '10b_executive_dashboard_scroll.png',
        'description': (
            'Scrolled executive view showing RM productivity comparison, Next Best Actions '
            'aggregated across all RMs, and the AI insight panels for each analytics widget.'
        ),
        'features': [
            'RM Cockpit: AUM leaderboard with days-since-last-contact and alert count per RM',
            'NBA Aggregation: top next-best-actions across the entire bank portfolio',
            'AI Insight panels: per-widget OCI Cohere analysis with actionable narrative',
            'Alert category breakdown by type and severity count',
            'Real-time data refresh from Oracle EXEC_AUM_MONTHLY and live customer tables',
        ],
    },

    # ── 11 COMPLIANCE ─────────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '11',
        'module_name': 'Compliance',
        'screenshot': '11_compliance.png',
        'description': (
            'Regulatory compliance monitoring module. Tracks KYC status, suitability '
            'assessment completeness, documentation expiry, and audit trail across '
            'all customer interactions recorded in Oracle ADB.'
        ),
        'features': [
            'KYC status overview: Verified / Pending / Expired breakdown across portfolio',
            'Suitability assessment compliance rate with drill-down by customer segment',
            'Document expiry alerts with 30/60/90-day advance warning',
            'Complete audit trail of all RM actions stored in AUDIT_LOG table',
            'Regulatory reporting summary exportable for compliance officer review',
        ],
    },

    # ── 12 ADMIN CONSOLE ─────────────────────────────────────────────────────
    {
        'type': 'module',
        'module_no': '12',
        'module_name': 'Admin Console',
        'screenshot': '12_admin_console.png',
        'description': (
            'System administration panel for managing users, RM assignments, system '
            'configuration, and monitoring Oracle ADB health. Accessible to manager '
            'role accounts only.'
        ),
        'features': [
            'User management: create/deactivate RM accounts and assign customer portfolios',
            'Oracle ADB connection health monitor with pool utilization metrics',
            'AI model configuration: manage OCI model endpoints and PAF agent settings',
            'System-wide settings: session timeout, cache TTL, alert thresholds',
            'Audit log viewer with filter by user, action type, and date range',
        ],
    },
    {
        'type': 'module',
        'module_no': '12',
        'module_name': 'Admin Console — Configuration',
        'screenshot': '12b_admin_console_detail.png',
        'description': (
            'Admin console detail view showing RM-to-customer assignment management '
            'and Oracle ADB system monitoring panels.'
        ),
        'features': [
            'RM–Customer assignment table with bulk reassignment capability',
            'Oracle ADB pool status with active connections and query performance',
            'PAF Agent activity log: recent agent invocations and response times',
            'Vector embedding status: customer profiles with up-to-date embeddings',
            'Database migration runner for schema updates (accessible to admin only)',
        ],
    },

    # ── ARCHITECTURE ─────────────────────────────────────────────────────────
    {
        'type': 'architecture',
        'title': 'Technical Architecture',
        'layers': [
            {
                'label': 'Frontend (Browser)',
                'color': C_CYAN,
                'items': [
                    'Single-Page Application (HTML/CSS/JavaScript)',
                    'Responsive dark-theme UI with Oracle-aligned design language',
                    'apiFetch() universal REST client with JWT auth headers',
                    'Chart.js for AUM trend, health distribution, CI radar visualizations',
                ],
            },
            {
                'label': 'Backend API (Node.js / Express)',
                'color': C_ACCENT,
                'items': [
                    'RESTful API: /api/auth, /api/customers, /api/executive, /api/copilot',
                    'JWT middleware for authentication + role-based authorization',
                    'Oracle AI Agent factory routes (PAF): maturity, reco, campaign, alerts',
                    'asyncHandler pattern with centralized error handling',
                ],
            },
            {
                'label': 'Oracle Database 26ai',
                'color': C_GOLD,
                'items': [
                    'Oracle Autonomous Database 26ai with Vector Search (VECTOR columns)',
                    'Tables: CUSTOMERS, PRODUCTS, HOLDINGS, ALERTS, RM_APPOINTMENTS, AUDIT_LOG',
                    'AI Insights cache: CUSTOMER_AI_INSIGHTS, EXEC_AUM_MONTHLY',
                    'Oracle SELECT AI: natural language SQL via AI Copilot',
                ],
            },
            {
                'label': 'OCI Generative AI (Cohere Command R+)',
                'color': C_GREEN,
                'items': [
                    'LLM: cohere.command-r-plus via OCI Generative AI Service',
                    'Oracle Private Agent Factory (PAF): tool-calling agents for RM workflows',
                    'RAG pipeline: vector embeddings → semantic search → context-augmented generation',
                    'Rule-based fallback when LLM unavailable (zero downtime)',
                ],
            },
        ],
    },

    # ── CLOSING ───────────────────────────────────────────────────────────────
    {
        'type': 'closing',
        'title': 'Intelligence RM Platform',
        'subtitle': 'Bank Danamon POC · 2026',
        'bullets': [
            '12 integrated modules across RM, Manager, and Admin roles',
            'Oracle AI · Oracle ADB 26ai · OCI Generative AI (Cohere Command R+)',
            'Oracle Private Agent Factory (PAF) · Oracle Select AI · Vector RAG',
            'Full audit trail · 24h AI insight cache · Rule-based fallback',
        ],
        'footer': 'Oracle AI · Bank Danamon Indonesia · Proof of Concept 2026',
    },
]


# ═══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_color=None, radius=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line)
    elif not line:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h, font_size=12, bold=False,
                 color=C_WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txBox


def set_slide_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_screenshot(slide, filename, l, t, w, h):
    img_path = os.path.join(SS_DIR, filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder if screenshot missing
        ph = add_rect(slide, l, t, w, h, fill=C_DIVIDER, line=0.5, line_color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ═══════════════════════════════════════════════════════════════════════════════

def build_cover(prs, data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # Left dark panel
    add_rect(slide, 0, 0, 5.8, 7.5, fill=C_CARD)

    # Oracle badge
    add_rect(slide, 0.55, 0.70, 1.20, 0.32, fill=C_ACCENT)
    add_text_box(slide, 'ORACLE', 0.55, 0.70, 1.20, 0.32,
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, 'AI · Bank Danamon POC', 1.82, 0.70, 3.5, 0.32,
                 font_size=9, color=C_LIGHT, align=PP_ALIGN.LEFT)

    # Big title
    add_text_box(slide, 'Intelligence', 0.55, 1.30, 5.0, 1.0,
                 font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, 'RM Platform', 0.55, 2.10, 5.0, 0.90,
                 font_size=44, bold=True, color=C_CYAN, align=PP_ALIGN.LEFT)

    # Subtitle
    add_text_box(slide, data['subtitle'], 0.55, 3.15, 5.0, 0.40,
                 font_size=14, bold=False, color=C_LIGHT)

    # Horizontal rule
    add_rect(slide, 0.55, 3.65, 4.8, 0.03, fill=C_ACCENT)

    # Tagline
    add_text_box(slide, data['tagline'], 0.55, 3.82, 5.0, 0.5,
                 font_size=9.5, color=C_LIGHT, italic=True)

    # Badges
    badge_x = 0.55
    for badge in data['badges']:
        add_rect(slide, badge_x, 4.55, 1.18, 0.30, fill=RGBColor(0x1A, 0x2A, 0x44),
                 line=0.5, line_color=C_CYAN)
        add_text_box(slide, badge, badge_x, 4.55, 1.18, 0.30,
                     font_size=7.5, color=C_CYAN, align=PP_ALIGN.CENTER)
        badge_x += 1.26

    # Year badge bottom-left
    add_text_box(slide, '2026', 0.55, 6.90, 1.0, 0.40,
                 font_size=11, color=RGBColor(0x40, 0x50, 0x70))

    # Right panel — decorative grid
    for row in range(12):
        for col in range(8):
            x = 6.0 + col * 0.92
            y = 0.3 + row * 0.59
            opacity_fill = RGBColor(0x15, 0x20, 0x35)
            add_rect(slide, x, y, 0.78, 0.45, fill=opacity_fill, line=0.3,
                     line_color=RGBColor(0x1E, 0x2D, 0x4A))

    # Big text on right
    add_text_box(slide, 'POC', 7.5, 2.8, 4.0, 2.0,
                 font_size=110, bold=True, color=RGBColor(0x18, 0x26, 0x3E),
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, '2026', 7.5, 5.5, 4.0, 1.0,
                 font_size=48, bold=True, color=RGBColor(0x1A, 0x2D, 0x48),
                 align=PP_ALIGN.CENTER)

    return slide


def build_agenda(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # Header bar
    add_rect(slide, 0, 0, 13.33, 1.1, fill=C_CARD)
    add_rect(slide, 0, 1.08, 13.33, 0.04, fill=C_ACCENT)
    add_text_box(slide, data['title'], 0.45, 0.20, 8.0, 0.70,
                 font_size=26, bold=True, color=C_WHITE)
    add_text_box(slide, 'Bank Danamon Intelligence RM Platform · POC 2026',
                 0.45, 0.70, 8.0, 0.38, font_size=10, color=C_LIGHT)

    # Two-column layout
    items = data['items']
    half = (len(items) + 1) // 2
    left_items  = items[:half]
    right_items = items[half:]

    for col_idx, col_items in enumerate([left_items, right_items]):
        col_x = 0.35 + col_idx * 6.55
        for i, (num, title, desc) in enumerate(col_items):
            y = 1.30 + i * 0.81
            # Number badge
            add_rect(slide, col_x, y + 0.04, 0.42, 0.35,
                     fill=C_ACCENT)
            add_text_box(slide, num, col_x, y + 0.04, 0.42, 0.35,
                         font_size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            # Card bg
            add_rect(slide, col_x + 0.50, y, 5.80, 0.67,
                     fill=C_CARD, line=0.5, line_color=C_DIVIDER)
            # Title
            add_text_box(slide, title, col_x + 0.62, y + 0.04, 5.50, 0.30,
                         font_size=11, bold=True, color=C_WHITE)
            # Desc
            add_text_box(slide, desc, col_x + 0.62, y + 0.33, 5.50, 0.28,
                         font_size=8.5, color=C_LIGHT)

    return slide


def build_module_slide(prs, data, screenshot_key='screenshot'):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # ── Header ────────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.33, 1.05, fill=C_CARD)
    add_rect(slide, 0, 1.03, 13.33, 0.04, fill=C_ACCENT)

    # Module number badge
    add_rect(slide, 0.35, 0.18, 0.55, 0.68, fill=C_ACCENT)
    add_text_box(slide, data['module_no'], 0.35, 0.18, 0.55, 0.68,
                 font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, data['module_name'], 1.02, 0.12, 9.0, 0.50,
                 font_size=20, bold=True, color=C_WHITE)
    add_text_box(slide, 'Bank Danamon Intelligence RM Platform', 1.02, 0.60,
                 9.0, 0.35, font_size=9, color=C_LIGHT)

    # Oracle badge (top-right)
    add_rect(slide, 11.90, 0.22, 1.08, 0.28, fill=C_ACCENT)
    add_text_box(slide, 'ORACLE AI', 11.90, 0.22, 1.08, 0.28,
                 font_size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Screenshot (left 60%) ─────────────────────────────────────────────────
    ss_filename = data.get(screenshot_key, '')
    add_screenshot(slide, ss_filename, 0.30, 1.22, 7.90, 5.10)

    # Light border around screenshot
    add_rect(slide, 0.30, 1.22, 7.90, 5.10, line=0.5, line_color=C_DIVIDER)

    # ── Right panel ───────────────────────────────────────────────────────────
    rx = 8.45
    # Description box
    add_rect(slide, rx, 1.22, 4.55, 1.70, fill=C_CARD, line=0.5, line_color=C_DIVIDER)
    add_text_box(slide, 'Overview', rx + 0.15, 1.28, 4.20, 0.28,
                 font_size=8, bold=True, color=C_CYAN)
    add_text_box(slide, data['description'], rx + 0.15, 1.56, 4.22, 1.28,
                 font_size=8.5, color=C_LIGHT)

    # Features
    features_y = 3.08
    add_text_box(slide, 'KEY FEATURES', rx, features_y, 4.55, 0.28,
                 font_size=8, bold=True, color=C_ACCENT)
    features_y += 0.28

    for feat in data.get('features', []):
        add_rect(slide, rx, features_y, 4.55, 0.58,
                 fill=C_CARD, line=0.3, line_color=C_DIVIDER)
        # bullet dot
        add_rect(slide, rx + 0.12, features_y + 0.22, 0.07, 0.07, fill=C_CYAN)
        add_text_box(slide, feat, rx + 0.28, features_y + 0.06, 4.15, 0.46,
                     font_size=8, color=C_LIGHT)
        features_y += 0.62

    # ── Footer ────────────────────────────────────────────────────────────────
    add_rect(slide, 0, 6.95, 13.33, 0.55, fill=C_CARD)
    add_text_box(slide, 'Oracle AI · Bank Danamon Intelligence RM Platform · POC 2026',
                 0.35, 7.00, 9.0, 0.40, font_size=8, color=RGBColor(0x40, 0x50, 0x70))
    add_text_box(slide, f'Module {data["module_no"]}',
                 12.20, 7.00, 1.0, 0.40, font_size=8, color=RGBColor(0x40, 0x50, 0x70),
                 align=PP_ALIGN.RIGHT)

    return slide


def build_architecture(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # Header
    add_rect(slide, 0, 0, 13.33, 1.05, fill=C_CARD)
    add_rect(slide, 0, 1.03, 13.33, 0.04, fill=C_ACCENT)
    add_text_box(slide, data['title'], 0.45, 0.12, 10.0, 0.55,
                 font_size=24, bold=True, color=C_WHITE)
    add_text_box(slide, 'Bank Danamon Intelligence RM Platform · Technical Stack',
                 0.45, 0.65, 10.0, 0.35, font_size=9, color=C_LIGHT)

    # Architecture diagram — 4 horizontal layer cards
    layer_colors = [C_CYAN, C_ACCENT, C_GOLD, C_GREEN]
    layers = data['layers']

    for i, layer in enumerate(layers):
        y = 1.22 + i * 1.47
        lc = layer_colors[i]

        # Left colour bar
        add_rect(slide, 0.30, y, 0.12, 1.25, fill=lc)

        # Card bg
        add_rect(slide, 0.45, y, 12.55, 1.25, fill=C_CARD, line=0.5, line_color=C_DIVIDER)

        # Layer label (vertical area)
        add_rect(slide, 0.45, y, 2.10, 1.25, fill=RGBColor(0x10, 0x18, 0x2A))
        add_text_box(slide, layer['label'], 0.55, y + 0.38, 1.90, 0.50,
                     font_size=9.5, bold=True, color=lc, align=PP_ALIGN.LEFT)

        # Items in 2-column grid
        col_items = layer['items']
        for j, item in enumerate(col_items):
            col = j % 2
            row = j // 2
            ix = 2.75 + col * 5.15
            iy = y + 0.12 + row * 0.55
            add_rect(slide, ix, iy, 4.95, 0.44,
                     fill=RGBColor(0x0D, 0x14, 0x22), line=0.3,
                     line_color=RGBColor(0x1E, 0x2D, 0x4A))
            add_rect(slide, ix + 0.08, iy + 0.17, 0.06, 0.06, fill=lc)
            add_text_box(slide, item, ix + 0.22, iy + 0.04, 4.60, 0.36,
                         font_size=8, color=C_LIGHT)

    # Footer
    add_rect(slide, 0, 6.95, 13.33, 0.55, fill=C_CARD)
    add_text_box(slide, 'Oracle AI · Bank Danamon Intelligence RM Platform · POC 2026',
                 0.35, 7.00, 9.0, 0.40, font_size=8, color=RGBColor(0x40, 0x50, 0x70))

    return slide


def build_closing(prs, data):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)

    # Background grid (decorative)
    for row in range(13):
        for col in range(15):
            x = col * 0.90
            y = row * 0.59
            add_rect(slide, x, y, 0.76, 0.45,
                     fill=RGBColor(0x10, 0x17, 0x28),
                     line=0.3, line_color=RGBColor(0x18, 0x22, 0x36))

    # Dark centre card
    add_rect(slide, 2.5, 1.20, 8.33, 5.10,
             fill=RGBColor(0x0C, 0x11, 0x1E),
             line=0.8, line_color=C_ACCENT)

    # Oracle badge
    add_rect(slide, 4.2, 1.52, 1.20, 0.32, fill=C_ACCENT)
    add_text_box(slide, 'ORACLE', 4.2, 1.52, 1.20, 0.32,
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 'AI · Bank Danamon POC', 5.48, 1.52, 3.5, 0.32,
                 font_size=9, color=C_LIGHT)

    add_text_box(slide, data['title'], 2.85, 2.00, 7.63, 0.90,
                 font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, data['subtitle'], 2.85, 2.88, 7.63, 0.45,
                 font_size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 3.5, 3.48, 6.33, 0.03, fill=C_ACCENT)

    # Bullets
    for i, bullet in enumerate(data['bullets']):
        by = 3.65 + i * 0.42
        add_rect(slide, 3.6, by + 0.12, 0.10, 0.10, fill=C_CYAN)
        add_text_box(slide, bullet, 3.82, by, 5.90, 0.38,
                     font_size=9.5, color=C_LIGHT)

    # Footer
    add_text_box(slide, data['footer'], 2.85, 5.95, 7.63, 0.35,
                 font_size=8.5, color=RGBColor(0x40, 0x50, 0x70),
                 align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for data in SLIDES:
        t = data['type']
        if t == 'cover':
            build_cover(prs, data)
            print('[✓] Cover slide')
        elif t == 'agenda':
            build_agenda(prs, data)
            print('[✓] Agenda slide')
        elif t in ('module', 'module_b'):
            build_module_slide(prs, data)
            print(f'[✓] Module {data["module_no"]} — {data["module_name"]}')
        elif t == 'architecture':
            build_architecture(prs, data)
            print('[✓] Architecture slide')
        elif t == 'closing':
            build_closing(prs, data)
            print('[✓] Closing slide')

    prs.save(OUT_PATH)
    print(f'\n[✓] Saved: {OUT_PATH}')
    print(f'[✓] Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
