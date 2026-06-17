"""
Add feature-evaluation slides and annotation callouts to the IRM documentation PPTX.
Evaluates 3 acceptance criteria against current IRM capabilities.
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX_PATH = r"C:\Users\deny\Desktop\Project-Danamon-RM\docs\Danamon_RM_Platform_Documentation.pptx"
OUT_PATH  = PPTX_PATH   # overwrite in-place

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x12, 0x1F)
C_CARD    = RGBColor(0x13, 0x1C, 0x30)
C_ACCENT  = RGBColor(0xFF, 0x4E, 0x00)
C_CYAN    = RGBColor(0x00, 0xCC, 0xFF)
C_GOLD    = RGBColor(0xFF, 0xB8, 0x30)
C_GREEN   = RGBColor(0x00, 0xD4, 0x7E)
C_RED     = RGBColor(0xFF, 0x4F, 0x4F)
C_AMBER   = RGBColor(0xFF, 0xB8, 0x30)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xB0, 0xBE, 0xD8)
C_DIVIDER = RGBColor(0x1E, 0x2D, 0x4A)
C_DARK2   = RGBColor(0x0A, 0x0F, 0x1A)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line_c=None, line_w=0.5):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    sp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.line.fill.background()
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_c:
        sp.line.color.rgb = line_c; sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    return sp

def tb(slide, text, l, t, w, h, size=10, bold=False, color=C_WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = 'Calibri'
    return bx

def header(slide, title, subtitle, module_no=None, color=C_ACCENT):
    rect(slide, 0, 0, 13.33, 1.05, fill=C_CARD)
    if module_no:
        rect(slide, 0.35, 0.18, 0.55, 0.68, fill=color)
        tb(slide, module_no, 0.35, 0.18, 0.55, 0.68,
           size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tx = 1.02
    else:
        tx = 0.45
    tb(slide, title, tx, 0.12, 10.0, 0.50, size=20, bold=True, color=C_WHITE)
    tb(slide, subtitle, tx, 0.62, 10.0, 0.35, size=9, color=C_LIGHT)
    rect(slide, 11.90, 0.22, 1.08, 0.28, fill=color)
    tb(slide, 'EVAL', 11.90, 0.22, 1.08, 0.28,
       size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def footer(slide, txt='Oracle AI · Bank Danamon Intelligence RM Platform · Feature Evaluation'):
    rect(slide, 0, 6.95, 13.33, 0.55, fill=C_CARD)
    tb(slide, txt, 0.35, 7.00, 11.0, 0.40, size=8,
       color=RGBColor(0x40, 0x50, 0x70))

def badge(slide, label, l, t, w=1.05, h=0.28, bg=C_GREEN, fg=C_WHITE):
    rect(slide, l, t, w, h, fill=bg)
    tb(slide, label, l, t, w, h, size=7.5, bold=True, color=fg,
       align=PP_ALIGN.CENTER)

def status_icon(slide, status, l, t):
    """status: 'met' | 'partial' | 'gap' """
    colors = {'met': C_GREEN, 'partial': C_AMBER, 'gap': C_RED}
    icons  = {'met': 'MET', 'partial': 'PARTIAL', 'gap': 'GAP'}
    bg = colors.get(status, C_LIGHT)
    lb = icons.get(status, status.upper())
    rect(slide, l, t, 0.72, 0.24, fill=bg)
    tb(slide, lb, l, t, 0.72, 0.24, size=7, bold=True,
       color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Annotation callout on an existing slide ───────────────────────────────────

def add_callout(slide, status, heading, lines, l=8.55, t=1.22, w=4.55):
    """
    Adds an evaluation callout panel on the RIGHT side of an existing module slide.
    status: 'MET' / 'PARTIAL' / 'GAP'
    """
    colors = {'MET': C_GREEN, 'PARTIAL': C_AMBER, 'GAP': C_RED}
    bg = colors.get(status, C_LIGHT)
    line_h = 0.52
    total_h = 0.50 + len(lines) * line_h + 0.15

    # Header bar
    rect(slide, l, t, w, 0.38, fill=bg)
    tb(slide, f'EVAL  {heading}', l + 0.10, t + 0.04, w - 0.15, 0.30,
       size=8.5, bold=True, color=C_WHITE)

    # Body
    rect(slide, l, t + 0.38, w, total_h - 0.38,
         fill=RGBColor(0x0A, 0x0F, 0x1A), line_c=bg, line_w=0.5)

    cur_y = t + 0.44
    for (s, txt) in lines:
        ic = {'v': C_GREEN, 'x': C_RED, '!': C_AMBER}.get(s, C_LIGHT)
        sym = {'v': '✔', 'x': '✘', '!': '⚠'}.get(s, '•')
        rect(slide, l + 0.10, cur_y + 0.08, 0.18, 0.18, fill=ic)
        tb(slide, sym, l + 0.10, cur_y + 0.06, 0.18, 0.22,
           size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(slide, txt, l + 0.34, cur_y + 0.04, w - 0.44, 0.44,
           size=7.5, color=C_LIGHT)
        cur_y += line_h

    return total_h

# ═══════════════════════════════════════════════════════════════════════════════
# Build evaluation summary slide
# ═══════════════════════════════════════════════════════════════════════════════

def build_eval_summary(prs):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_bg(slide)
    header(slide,
           'Feature Evaluation Summary',
           'IRM Capability Assessment against Acceptance Criteria · 3 Goals',
           color=C_CYAN)
    footer(slide)

    # ── Legend row ────────────────────────────────────────────────
    lx = 0.45
    for lbl, clr in [('MET', C_GREEN), ('PARTIAL', C_AMBER), ('GAP — Not Available', C_RED)]:
        rect(slide, lx, 1.14, 0.12, 0.12, fill=clr)
        tb(slide, lbl, lx + 0.18, 1.10, 1.40, 0.22, size=8, color=C_LIGHT)
        lx += 1.75

    GOALS = [
        {
            'no': '01',
            'title': 'Maturity Reminder with Action Plan',
            'module': 'Module 05 — Maturity Reminder',
            'goal': 'Automated reminders when a product matures, with a detailed action plan including talking points.',
            'criteria': [
                ('met',     'Detects maturities within 30/60/90-day horizon from Oracle ADB'),
                ('met',     'AI generates reinvestment recommendation action plan per customer'),
                ('met',     'Action Plan Follow-up button, share via Email / WhatsApp'),
                ('partial', 'Talking points are reinvestment-focused; not structured per customer dialogue script'),
                ('gap',     'System is ON-DEMAND only — no automated push/email notification scheduler'),
                ('gap',     'No supporting documents or product brochures attached to reminders'),
            ],
            'verdict': 'PARTIAL',
        },
        {
            'no': '02',
            'title': 'Product Profile Matching & Recommendation',
            'module': 'Module 06 — Product Recommendations',
            'goal': 'List of products matching the customer\'s profile with justification, details, risks, and goals alignment.',
            'criteria': [
                ('met',     'Profile Matching: risk tolerance, portfolio gaps, transaction history'),
                ('met',     'Recommendation List: ranked products with fit score per customer'),
                ('met',     'Recommendation Justification: multi-factor AI scoring + Bahasa rationale'),
                ('partial', 'Product Details: fit score & rationale shown; explicit risk/benefit/tenor per product not confirmed'),
                ('gap',     'Recommendations are batch/portfolio-wide — no per-customer on-demand lookup'),
                ('gap',     'No product detail sheet or brochure link shown alongside each recommendation'),
            ],
            'verdict': 'PARTIAL',
        },
        {
            'no': '03',
            'title': 'Portfolio Alerts with Actionable Insights',
            'module': 'Module 08 — Portfolio Alerts',
            'goal': 'Alerts for significant portfolio changes with reason, affected assets, and actionable recommendations.',
            'criteria': [
                ('met',     'Alert triggers: Churn Risk, Idle Funds, KYC Expiry, Maturity, AUM Drop'),
                ('met',     'Per-alert card: affected investment, reason, severity badge'),
                ('met',     'AI-generated remediation steps ranked by impact; link to Customer 360'),
                ('partial', 'Suggested actions include discuss/review; direct rebalancing workflow not present'),
                ('gap',     'Alert parameters NOT configurable by Product Team — criteria are system-hardcoded'),
                ('gap',     'Market-event-driven alerts (e.g. IHSG drop, rate hike) not implemented'),
            ],
            'verdict': 'PARTIAL',
        },
    ]

    goal_colors = [C_CYAN, C_GOLD, C_ACCENT]
    verdict_colors = {'MET': C_GREEN, 'PARTIAL': C_AMBER, 'GAP': C_RED}

    for gi, g in enumerate(GOALS):
        gy = 1.42 + gi * 1.80
        gc = goal_colors[gi]
        vc = verdict_colors[g['verdict']]

        # Left badge
        rect(slide, 0.30, gy, 0.42, 1.55, fill=gc)
        tb(slide, g['no'], 0.30, gy + 0.55, 0.42, 0.50,
           size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Main card
        rect(slide, 0.75, gy, 12.25, 1.55, fill=C_CARD, line_c=C_DIVIDER)

        # Verdict badge top-right
        rect(slide, 11.60, gy + 0.14, 1.20, 0.28, fill=vc)
        tb(slide, g['verdict'], 11.60, gy + 0.14, 1.20, 0.28,
           size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Title + module
        tb(slide, g['title'], 0.92, gy + 0.08, 10.0, 0.32,
           size=12, bold=True, color=C_WHITE)
        tb(slide, f"{g['module']}  ·  {g['goal']}", 0.92, gy + 0.38, 10.50, 0.32,
           size=8, color=C_LIGHT, italic=True)

        # Criteria dots: 3 per row
        for ci, (st, txt) in enumerate(g['criteria']):
            col = ci % 3
            row = ci // 3
            cx = 0.92 + col * 4.08
            cy = gy + 0.76 + row * 0.38
            sc = {'met': C_GREEN, 'partial': C_AMBER, 'gap': C_RED}[st]
            sym = {'met': '✔', 'partial': '⚠', 'gap': '✘'}[st]
            rect(slide, cx, cy + 0.06, 0.16, 0.16, fill=sc)
            tb(slide, sym, cx, cy + 0.04, 0.16, 0.22,
               size=6.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            tb(slide, txt, cx + 0.22, cy, 3.72, 0.34, size=7.5, color=C_LIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# Build 3 detailed evaluation slides (one per goal)
# ═══════════════════════════════════════════════════════════════════════════════

def build_eval_detail(prs, goal_no, title, module_ref, goal_txt, benefit_txt,
                      met_items, gap_items, reco_items, verdict, header_color):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_bg(slide)

    verdict_colors = {'MET': C_GREEN, 'PARTIAL': C_AMBER, 'GAP': C_RED}
    vc = verdict_colors.get(verdict, C_AMBER)

    # Header
    header(slide, f'Goal {goal_no}: {title}', f'{module_ref}  ·  Feature Gap Analysis',
           module_no=goal_no, color=header_color)
    footer(slide)

    # Verdict banner
    rect(slide, 0.30, 1.14, 12.73, 0.30, fill=vc)
    verdict_txt = {'MET': 'FULLY MET', 'PARTIAL': 'PARTIALLY MET — Gaps Identified',
                   'GAP': 'NOT MET — Feature Missing'}[verdict]
    tb(slide, f'OVERALL VERDICT:  {verdict_txt}', 0.45, 1.16, 12.0, 0.26,
       size=9, bold=True, color=C_WHITE)

    # ── Goal / Benefit / Criteria box (top-left) ───────────────────
    rect(slide, 0.30, 1.52, 6.10, 1.62, fill=C_CARD, line_c=C_DIVIDER)
    tb(slide, 'GOAL & BENEFIT', 0.45, 1.57, 5.80, 0.22, size=7.5, bold=True, color=header_color)
    tb(slide, f'Goal: {goal_txt}', 0.45, 1.80, 5.82, 0.38, size=8.5, color=C_WHITE)
    tb(slide, f'Benefit: {benefit_txt}', 0.45, 2.19, 5.82, 0.38, size=8, color=C_LIGHT, italic=True)
    tb(slide, f'Module: {module_ref}', 0.45, 2.59, 5.82, 0.22, size=8, color=header_color)

    # ── What IRM provides (MET items) ─────────────────────────────
    rect(slide, 6.55, 1.52, 6.48, 1.62, fill=C_CARD, line_c=C_DIVIDER)
    tb(slide, 'WHAT IRM CURRENTLY PROVIDES', 6.70, 1.57, 6.20, 0.22,
       size=7.5, bold=True, color=C_GREEN)
    for i, item in enumerate(met_items[:4]):
        iy = 1.82 + i * 0.32
        rect(slide, 6.70, iy + 0.07, 0.14, 0.14, fill=C_GREEN)
        tb(slide, '✔', 6.70, iy + 0.05, 0.14, 0.20, size=6.5, bold=True,
           color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(slide, item, 6.90, iy, 6.00, 0.28, size=8.5, color=C_LIGHT)

    # ── Gap analysis (2 columns) ───────────────────────────────────
    tb(slide, 'ACCEPTANCE CRITERIA GAPS', 0.30, 3.22, 12.73, 0.22,
       size=8, bold=True, color=C_RED)

    for i, (crit, analysis, status) in enumerate(gap_items):
        col = i % 2
        row = i // 2
        gx = 0.30 + col * 6.46
        gy = 3.48 + row * 1.05
        sc = {'gap': C_RED, 'partial': C_AMBER, 'met': C_GREEN}.get(status, C_AMBER)
        lbl = {'gap': 'GAP', 'partial': 'PARTIAL', 'met': 'MET'}.get(status, 'PARTIAL')

        rect(slide, gx, gy, 6.28, 0.95, fill=RGBColor(0x0A, 0x0F, 0x1A), line_c=sc, line_w=0.5)
        rect(slide, gx, gy, 0.75, 0.28, fill=sc)
        tb(slide, lbl, gx, gy, 0.75, 0.28, size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(slide, crit, gx + 0.82, gy + 0.04, 5.35, 0.24, size=8, bold=True, color=C_WHITE)
        tb(slide, analysis, gx + 0.12, gy + 0.32, 6.06, 0.58, size=7.8, color=C_LIGHT)

    # ── Recommendations ───────────────────────────────────────────
    ry = 3.48 + ((len(gap_items) + 1) // 2) * 1.05 + 0.10
    if ry > 5.80:
        ry = 5.80
    rect(slide, 0.30, ry, 12.73, 0.22, fill=C_CARD)
    tb(slide, 'RECOMMENDATIONS FOR IMPROVEMENT', 0.45, ry + 0.02, 12.0, 0.18,
       size=8, bold=True, color=C_GOLD)

    for ri, reco in enumerate(reco_items):
        rx2 = 0.30 + (ri % 3) * 4.27
        ry2 = ry + 0.28 + (ri // 3) * 0.42
        rect(slide, rx2, ry2, 4.15, 0.34, fill=C_CARD, line_c=C_DIVIDER)
        rect(slide, rx2, ry2, 0.30, 0.34, fill=C_GOLD)
        tb(slide, str(ri + 1), rx2, ry2, 0.30, 0.34, size=8, bold=True,
           color=C_DARK2, align=PP_ALIGN.CENTER)
        tb(slide, reco, rx2 + 0.36, ry2 + 0.04, 3.72, 0.30, size=8, color=C_LIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# Add annotation callout to existing slide (non-destructive — adds shapes)
# ═══════════════════════════════════════════════════════════════════════════════

CALLOUT_SPECS = {
    # slide_index (0-based) → (status, heading, lines)
    14: ('PARTIAL', 'Goal 1: Maturity Reminder', [        # slide 15
        ('v', 'Maturity detection within 30/60/90-day horizon'),
        ('v', 'AI reinvestment action plan generated per customer'),
        ('!', 'Talking points are generic; no per-customer dialogue script'),
        ('x', 'ON-DEMAND only — no automated push/email scheduler'),
        ('x', 'No product documents/brochures attached to reminder'),
    ]),
    15: ('PARTIAL', 'Goal 1: Maturity Result', [           # slide 16
        ('v', 'Action Plan Follow-up CTA button available'),
        ('v', 'Share via Email / WhatsApp for customer outreach'),
        ('!', 'Action plan lacks structured talking-point format'),
        ('x', 'No supporting data/document attachment per reminder'),
    ]),
    16: ('PARTIAL', 'Goal 2: Product Reco Input', [        # slide 17
        ('v', 'Risk profile, portfolio gap, AUM potential used for matching'),
        ('v', 'Multi-factor scoring: risk alignment + urgency'),
        ('!', 'Recommendations are batch-wide, not per-customer on demand'),
        ('x', 'No product risk/benefit/tenor detail displayed per item'),
    ]),
    17: ('PARTIAL', 'Goal 2: Reco Results', [              # slide 18
        ('v', 'Ranked list with fit score + AI rationale in Bahasa'),
        ('!', 'Product details (risks, tenor, return) not explicitly shown'),
        ('x', 'No individual customer lookup — must run full portfolio scan'),
    ]),
    20: ('PARTIAL', 'Goal 3: Portfolio Alerts', [          # slide 21
        ('v', 'Alert categories: Churn, Idle Funds, KYC, Maturity, AUM Drop'),
        ('v', 'Severity: Critical / High / Medium / Low with badges'),
        ('x', 'Criteria are hardcoded — Product Team cannot configure params'),
        ('x', 'No market-event-driven alerts (IHSG drop, BI Rate change)'),
    ]),
    21: ('PARTIAL', 'Goal 3: Alert Detail', [              # slide 22
        ('v', 'AI root-cause explanation + ranked remediation steps'),
        ('v', 'One-click to Customer 360 for investigation'),
        ('!', 'Suggested actions are text-only — no direct rebalancing flow'),
        ('x', 'No "underperforming asset" sell/rebalance action trigger'),
    ]),
}


def annotate_existing_slides(prs):
    for idx, (status, heading, lines) in CALLOUT_SPECS.items():
        if idx >= len(prs.slides):
            print(f'  [skip] slide index {idx} out of range')
            continue
        slide = prs.slides[idx]
        color_map = {'PARTIAL': C_AMBER, 'MET': C_GREEN, 'GAP': C_RED}
        bar_color = color_map.get(status, C_AMBER)

        # Semi-transparent overlay strip on right panel to indicate evaluation
        rect(slide, 8.45, 1.15, 4.58, 0.26, fill=bar_color)
        tb(slide, f'EVALUATION  ·  {status}', 8.55, 1.17, 4.40, 0.22,
           size=8, bold=True, color=C_WHITE)

        # Lines in evaluation panel (replace right-side content area)
        panel_y = 1.44
        for (sym, txt) in lines:
            ic = {'v': C_GREEN, 'x': C_RED, '!': C_AMBER}.get(sym, C_LIGHT)
            icon_char = {'v': '✔', 'x': '✘', '!': '⚠'}.get(sym, '•')
            rect(slide, 8.55, panel_y + 0.07, 0.18, 0.18, fill=ic)
            tb(slide, icon_char, 8.55, panel_y + 0.05, 0.18, 0.22,
               size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            tb(slide, txt, 8.79, panel_y + 0.02, 4.10, 0.42, size=8, color=C_LIGHT)
            panel_y += 0.48

        print(f'  [annotated] slide {idx + 1}')


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(PPTX_PATH)
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print('\n[1] Annotating existing module slides...')
    annotate_existing_slides(prs)

    print('\n[2] Building evaluation slides...')

    # ── Eval Summary (insert after slide 2 = agenda) ──────────────
    summary_slide = build_eval_summary(prs)
    print('  [+] Evaluation Summary slide')

    # ── Goal 1 Detail ─────────────────────────────────────────────
    d1 = build_eval_detail(
        prs,
        goal_no='01',
        title='Maturity Reminder with Action Plan',
        module_ref='Module 05 — Maturity Reminder',
        goal_txt='Receive automated reminders when a product matures, with a suggested action plan.',
        benefit_txt='RM can immediately strategize a rearrangement strategy for customers.',
        met_items=[
            'Detects upcoming maturities within 30/60/90-day configurable horizon',
            'AI generates personalized reinvestment recommendation per customer',
            'Action Plan Follow-up CTA available after each analysis result',
            'Share via Email and WhatsApp for direct customer outreach',
        ],
        gap_items=[
            ('Action Plan Format',
             'The action plan provides reinvestment product suggestions but lacks a structured '
             'customer dialogue script with specific talking points (e.g., objection handling, '
             'yield comparison narrative, persuasion cues).',
             'partial'),
            ('Automated Push Notifications',
             'System is ON-DEMAND only: RM must manually trigger analysis. There is no scheduled '
             'push notification, email alert, or in-app reminder that fires automatically when '
             'a product reaches maturity date.',
             'gap'),
            ('Supporting Documents',
             'Reminders do not include attached product brochures, term-sheet PDFs, or factsheet '
             'links to support the conversation. RM must source these separately.',
             'gap'),
            ('Specific Criteria Scope',
             'Acceptance criterion: "Each reminder includes a detailed action plan with specific '
             'talking points and supporting data." Current output is AI-text only, without '
             'structured per-product documentation.',
             'partial'),
        ],
        reco_items=[
            'Add scheduled job (cron/Oracle DBMS_SCHEDULER) to push maturity alerts automatically',
            'Structure action plan output with labelled sections: Opening, Value Proposition, Objection Handling, Close',
            'Attach product factsheet URL or PDF link to each recommended reinvestment product',
            'Integrate notification channel: in-app badge + optional email digest',
            'Allow RM to set personal reminder horizon (e.g., alert me 14 days before maturity)',
        ],
        verdict='PARTIAL',
        header_color=C_CYAN,
    )
    print('  [+] Goal 01 detail slide')

    # ── Goal 2 Detail ─────────────────────────────────────────────
    d2 = build_eval_detail(
        prs,
        goal_no='02',
        title='Product Profile Matching & Recommendation',
        module_ref='Module 06 — Product Recommendations',
        goal_txt='See a list of products that fit the customer\'s profile with justification and product details.',
        benefit_txt='RM can confidently recommend the most suitable products for each customer.',
        met_items=[
            'Profile matching uses: risk tolerance, portfolio gaps, transaction history, market context',
            'Ranked recommendation list with fit score + AI rationale (Bahasa Indonesia)',
            'Multi-factor scoring: portfolio gap, risk alignment, AUM potential, urgency',
            'Share CTA + Action Plan Follow-up per recommended customer',
        ],
        gap_items=[
            ('Product Details per Recommendation',
             'Each recommended product should show: benefits, risks, tenor/duration, projected return, '
             'and alignment with customer goals. Current output shows product name + fit rationale but '
             'does not display a full product detail card.',
             'partial'),
            ('Per-Customer On-Demand Lookup',
             'Recommendations are generated as a batch for the entire RM portfolio. There is no '
             '"recommend products for THIS customer" feature accessible directly from the '
             'Customer 360 profile page.',
             'gap'),
            ('Product Catalog Detail Page',
             'No product brochure or detail sheet is linked or displayed alongside each recommendation. '
             'RM cannot drill into product terms, rates, or risk rating from within the platform.',
             'gap'),
            ('Customer Goal Alignment',
             'Criterion: "alignment with customer goals." The system uses risk profile and AUM; '
             'explicit customer financial goals (retirement, education, etc.) are not captured '
             'or used as a matching dimension.',
             'partial'),
        ],
        reco_items=[
            'Add Product Detail panel/modal: show tenor, return rate, risk rating, min invest, and goal tag',
            'Add "Get Recommendations" button on Customer 360 profile for per-customer on-demand reco',
            'Capture customer financial goals during onboarding and use as a reco filter dimension',
            'Link product factsheet PDF or bank product page URL alongside each recommendation card',
            'Show product comparison view (side-by-side) for top-3 recommended products',
        ],
        verdict='PARTIAL',
        header_color=C_GOLD,
    )
    print('  [+] Goal 02 detail slide')

    # ── Goal 3 Detail ─────────────────────────────────────────────
    d3 = build_eval_detail(
        prs,
        goal_no='03',
        title='Portfolio Alerts with Actionable Insights',
        module_ref='Module 08 — Portfolio Alerts',
        goal_txt='Receive alerts for significant portfolio changes/issues with reason, affected assets, and actionable recommendations.',
        benefit_txt='RM can promptly address significant changes and take appropriate action to manage portfolios effectively.',
        met_items=[
            'Alert categories: Churn Risk, Idle Funds, KYC Expiry, Maturity Approaching, AUM Drop',
            'Per-alert card: affected investment name, risk type, severity badge (Critical–Low)',
            'AI-generated root-cause explanation and ranked remediation steps',
            'One-click Customer 360 link for immediate in-depth investigation',
        ],
        gap_items=[
            ('Product-Team Configurable Parameters',
             'Criterion: "alerts triggered based on predefined criteria that can be set by Product Team." '
             'Current alert triggers (churn threshold, AUM drop %, idle days) are hardcoded in the system. '
             'There is no configuration UI for Product Team to set or adjust alert parameters.',
             'gap'),
            ('Market-Event-Driven Alerts',
             'No alerts are triggered by external market events (e.g., IHSG drops >3%, BI Rate change, '
             'USD/IDR volatility). Alerts are purely portfolio-level, not market-condition-aware.',
             'gap'),
            ('Direct Rebalancing Workflow',
             'Criterion: "actionable insights such as rebalancing assets or selling underperforming '
             'investments." Current alerts provide text recommendations only; no in-platform '
             'workflow exists to directly initiate a rebalancing or sell action.',
             'partial'),
            ('Underperforming Asset Detection',
             'No alert type specifically monitors for underperforming assets against benchmark '
             'returns. Current alerts focus on lifecycle events (maturity, KYC), not performance.',
             'gap'),
        ],
        reco_items=[
            'Build Admin/Product Team configuration panel for alert thresholds (AUM drop %, idle days, etc.)',
            'Integrate market-event triggers: connect to live market data feed → fire alerts on IHSG/rate events',
            'Add "Underperforming Asset" alert type comparing holding returns vs. benchmark/peer',
            'Add direct action buttons in alert: "Schedule Discussion," "Initiate Rebalancing," "Create Follow-up Task"',
            'Enable RM to subscribe/unsubscribe to specific alert categories per customer segment',
        ],
        verdict='PARTIAL',
        header_color=C_ACCENT,
    )
    print('  [+] Goal 03 detail slide')

    # ── Reorder slides: inject evaluation slides after slide index 2 ──
    # Current order after adding: ... summary, d1, d2, d3 appended at end
    # We want them at positions 3, 4, 5, 6 (0-indexed: 2, 3, 4, 5)
    # python-pptx doesn't support slide reordering natively; we manipulate XML
    xml_slides = prs.slides._sldIdLst
    all_ids    = list(xml_slides)

    # The 4 new slides are the last 4 elements
    new_ids = all_ids[-4:]   # [summary, d1, d2, d3]
    old_ids = all_ids[:-4]   # original 31 slides

    # Insert after position 2 (0-indexed = after agenda slide)
    reordered = old_ids[:2] + new_ids + old_ids[2:]

    # Rebuild the sldIdLst
    for el in list(xml_slides):
        xml_slides.remove(el)
    for el in reordered:
        xml_slides.append(el)

    print('\n[3] Saving presentation...')
    prs.save(OUT_PATH)
    print(f'[✓] Saved: {OUT_PATH}')
    print(f'[✓] Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
